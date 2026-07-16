import sys, io, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = (r'c:\Users\prathmesh.bargal\Desktop\Prathmesh\Case Study'
       r'\Case study\ds_candidate_assignment\Spendy_Loyalty_Presentation_Final.pptx')
DST = (r'c:\Users\prathmesh.bargal\Desktop\Prathmesh\Case Study'
       r'\Case study\ds_candidate_assignment\Spendy_Loyalty_Presentation_Premium.pptx')
shutil.copy2(SRC, DST)
prs = Presentation(DST)

GOLD  = RGBColor(0xC0, 0x98, 0x40)
NAVY  = RGBColor(0x0B, 0x2A, 0x4A)
STEEL = RGBColor(0x2E, 0x75, 0xB6)
MBLUE = RGBColor(0x2D, 0x6A, 0x9F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY  = RGBColor(0x95, 0xA5, 0xA6)
LGREY = RGBColor(0xF2, 0xF5, 0xF8)
DARK  = RGBColor(0x2C, 0x3E, 0x50)
SLATE = RGBColor(0x5E, 0x70, 0x80)
SBLUE = RGBColor(0xB0, 0xC8, 0xE4)

def emu(x): return int(x * 914400)

def add_rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(0.5)):
    s = slide.shapes.add_shape(1, emu(l), emu(t), emu(w), emu(h))
    s.line.width = Pt(0)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    return s

def add_tb(slide, l, t, w, h, text, sz, bold=False, color=WHITE,
           align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(emu(l), emu(t), emu(w), emu(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = sz; r.font.bold = bold; r.font.color.rgb = color
    return tb

def fc(shape, color):
    if not shape.has_text_frame: return
    for p in shape.text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb = color

# ── SLIDE 1: number badges 01-05 → GOLD ────────────────────
sl1 = prs.slides[0]
for s in sl1.shapes:
    if s.name in {'TextBox 11','TextBox 16','TextBox 21','TextBox 26','TextBox 31'}:
        fc(s, GOLD)
print("S1: badges GOLD")

# ── SLIDE 2: shrink body text + add 75% hero box ────────────
sl2 = prs.slides[1]
for s in sl2.shapes:
    if s.name in ('TextBox 8','TextBox 12','TextBox 16'):
        s.width = emu(8.2)
# Navy hero column on right
add_rect(sl2, 9.05, 1.35, 3.95, 5.5, fill=NAVY)
add_tb(sl2, 9.05, 1.65, 3.95, 1.0,  '75%',
       Pt(64), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_rect(sl2, 9.35, 2.80, 3.35, 0.025, fill=GOLD)
add_tb(sl2, 9.05, 2.88, 3.95, 0.32,
       'of users shop at exactly one merchant',
       Pt(10), color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
add_tb(sl2, 9.05, 3.28, 3.95, 0.28,
       'Merchant loyalty  ≠  Platform loyalty',
       Pt(9.5), bold=True, color=SBLUE, align=PP_ALIGN.CENTER)
add_tb(sl2, 9.05, 3.70, 3.95, 0.65,
       'If that merchant leaves Spendy,\nthe customer disappears entirely.',
       Pt(9), color=SBLUE, align=PP_ALIGN.CENTER, wrap=True)
print("S2: hero 75% box added")

# ── SLIDE 5: teal loyalty component column → MBLUE ──────────
sl5 = prs.slides[4]
for s in sl5.shapes:
    if not s.has_text_frame: continue
    for p in s.text_frame.paragraphs:
        for r in p.runs:
            try:
                if str(r.font.color.rgb) == '3D6B7A': r.font.color.rgb = MBLUE
            except: pass
print("S5: teal → MBLUE")

# ── SLIDE 7: recall 0.45 → GOLD; KEY INSIGHT gold border ────
sl7 = prs.slides[6]
for s in sl7.shapes:
    if s.name == 'TextBox 95': fc(s, GOLD)
    elif s.name == 'Rectangle 77':
        s.line.color.rgb = GOLD; s.line.width = Pt(2)
print("S7: recall GOLD, KEY INSIGHT border")

# ── SLIDE 8: LOW stat → GREY, MEDIUM stat → STEEL ───────────
sl8 = prs.slides[7]
for s in sl8.shapes:
    if s.name == 'TextBox 12': fc(s, GREY)
    elif s.name == 'TextBox 24': fc(s, STEEL)
print("S8: LOW grey, MED steel")

# ── SLIDE 9: 4.6% and 83.9% → GOLD ─────────────────────────
sl9 = prs.slides[8]
for s in sl9.shapes:
    if s.name == 'TextBox 12': fc(s, GOLD)
    elif s.name == 'TextBox 26': fc(s, GOLD)
print("S9: 4.6% and 83.9% GOLD")

# ── SLIDE 10: rebuild as 3+2 card grid ──────────────────────
sl10 = prs.slides[9]
KEEP = {'Rectangle 1','Rectangle 2','TextBox 3','TextBox 4'}
sp10 = sl10.shapes._spTree
for s in list(sl10.shapes):
    if s.name not in KEEP:
        sp10.remove(s._element)

HDR = [NAVY, MBLUE, STEEL, NAVY, MBLUE]
TITLES = [
    'HEURISTIC THRESHOLDS',
    'RECALL AT 45%',
    'NARROW FEATURE SET',
    '7-MONTH DATA WINDOW',
    'FIXED SEGMENTATION',
]
LIMITS = [
    'Loyalty thresholds (2/month, ¥10K, 2 shops, 6 months) set by data inspection, not optimised against outcome data.',
    'Model misses 55% of loyal users at the 0.5 threshold. Class imbalance (1 in 31) limits raw recall.',
    'Only 5 features from 3 raw columns. Cannot capture product affinity, app engagement, or channel response.',
    'H2 2020 data only. Post-COVID patterns and seasonality may not generalise across a full calendar year.',
    '33%/66% probability cutoffs are heuristic, not data-driven. Optimal thresholds depend on campaign ROI.',
]
FIXES = [
    'A/B test threshold variants; back-solve from LTV or reward ROI to find the optimal combination.',
    'Test SMOTE + Gradient Boosting; lower threshold to 0.3 to prioritise recall over precision.',
    'Add product category diversity, app engagement metrics, and marketing-channel response indicators.',
    'Extend training to 2+ years; cross-validate across cohorts and seasonal periods.',
    'Set cutoffs via decision theory: find the probability threshold where expected ROI per action is maximised.',
]

CW3 = 4.09; CW2 = 6.35; GAP = 0.13; CH = 2.65
Y1 = 1.25; Y2 = 4.10; X0 = 0.3
MGREY = RGBColor(0xCC, 0xD8, 0xE4)

for i in range(3):
    x = X0 + i * (CW3 + GAP)
    y = Y1; cw = CW3
    add_rect(sl10, x, y, cw, CH, fill=LGREY, border=MGREY)
    add_rect(sl10, x, y, cw, 0.42, fill=HDR[i])
    add_tb(sl10, x+0.1, y+0.05, cw-0.15, 0.32, TITLES[i], Pt(9), bold=True, color=WHITE)
    add_tb(sl10, x+0.12, y+0.52, cw-0.22, 0.20, 'CURRENT LIMITATION', Pt(7.5), bold=True, color=SLATE)
    add_tb(sl10, x+0.12, y+0.74, cw-0.22, 0.88, LIMITS[i], Pt(8.5), color=DARK, wrap=True)
    add_rect(sl10, x+0.12, y+1.63, cw-0.24, 0.02, fill=GOLD)
    add_tb(sl10, x+0.12, y+1.68, 0.6, 0.22, '→ FIX:', Pt(8), bold=True, color=GOLD)
    add_tb(sl10, x+0.12, y+1.94, cw-0.22, 0.68, FIXES[i], Pt(8.5), color=DARK, wrap=True)

for i in range(2):
    x = X0 + i * (CW2 + GAP)
    y = Y2; cw = CW2
    add_rect(sl10, x, y, cw, CH, fill=LGREY, border=MGREY)
    add_rect(sl10, x, y, cw, 0.42, fill=HDR[i+3])
    add_tb(sl10, x+0.1, y+0.05, cw-0.15, 0.32, TITLES[i+3], Pt(9), bold=True, color=WHITE)
    add_tb(sl10, x+0.12, y+0.52, cw-0.22, 0.20, 'CURRENT LIMITATION', Pt(7.5), bold=True, color=SLATE)
    add_tb(sl10, x+0.12, y+0.74, cw-0.22, 0.88, LIMITS[i+3], Pt(9.5), color=DARK, wrap=True)
    add_rect(sl10, x+0.12, y+1.63, cw-0.24, 0.02, fill=GOLD)
    add_tb(sl10, x+0.12, y+1.68, 0.6, 0.22, '→ FIX:', Pt(8), bold=True, color=GOLD)
    add_tb(sl10, x+0.12, y+1.94, cw-0.22, 0.68, FIXES[i+3], Pt(9.5), color=DARK, wrap=True)
print("S10: 3+2 grid built")

# ── SLIDE 11: gold rule + ①②③ → GOLD+larger + closing banner
sl11 = prs.slides[10]
add_rect(sl11, 0.3, 1.20, 12.73, 0.03, fill=GOLD)
for s in sl11.shapes:
    if s.name in ('TextBox 9','TextBox 13','TextBox 17'):
        for p in s.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(17); r.font.color.rgb = GOLD
add_rect(sl11, 0.3, 6.55, 12.73, 0.70, fill=NAVY)
add_tb(sl11, 0.3, 6.60, 12.73, 0.60,
       'Rigorous definition.     Month-2 scoring.     Three actionable segments.',
       Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
print("S11: gold rule + badges + closing banner")

# ── SLIDE 12: gold accent bar on RF column + CHOSEN → GOLD ──
sl12 = prs.slides[11]
add_rect(sl12, 3.54, 1.10, 0.07, 6.1, fill=GOLD)
for s in sl12.shapes:
    if s.name == 'TextBox 16': fc(s, GOLD)
print("S12: gold accent bar + CHOSEN badge GOLD")

# ── SLIDE 13: standardise pipeline node colors ───────────────
sl13 = prs.slides[12]
NODE_COLORS = {
    'Rectangle 9':  NAVY,
    'Rectangle 13': MBLUE,
    'Rectangle 17': STEEL,
    'Rectangle 21': MBLUE,   # was teal 3D6B7A
    'Rectangle 25': STEEL,   # was slate 5E7080
    'Rectangle 29': NAVY,    # was tan A09080
}
for s in sl13.shapes:
    if s.name in NODE_COLORS:
        s.fill.solid(); s.fill.fore_color.rgb = NODE_COLORS[s.name]
print("S13: node colors standardised")

prs.save(DST)
print(f"\nSaved: {DST}")
