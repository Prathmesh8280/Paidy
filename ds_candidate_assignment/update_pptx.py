import sys, io, os, html
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX = r'c:\Users\prathmesh.bargal\Desktop\Prathmesh\Case Study\Case study\ds_candidate_assignment\Spendy_Loyalty_Presentation_Final.pptx'
IMGS = r'c:\Users\prathmesh.bargal\Desktop\Prathmesh\Case Study\Case study\ds_candidate_assignment\chart_images'

prs = Presentation(PPTX)
W = prs.slide_width   # 13.33"
H = prs.slide_height  # 7.5"

# ── colour constants ──────────────────────────────────────────────────────────
NAVY  = RGBColor(0x0B, 0x2A, 0x4A)
STEEL = RGBColor(0x2E, 0x75, 0xB6)
TEAL  = RGBColor(0x3D, 0x6B, 0x7A)
BLUE  = RGBColor(0x2D, 0x6A, 0x9F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGREY = RGBColor(0xF5, 0xF5, 0xF5)
MGREY = RGBColor(0xCC, 0xCC, 0xCC)
DARK  = RGBColor(0x2C, 0x3E, 0x50)
AMBER = RGBColor(0xC0, 0x98, 0x40)

NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ── helpers ───────────────────────────────────────────────────────────────────

def emu(inches): return int(inches * 914400)

def add_rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(0.5)):
    shape = slide.shapes.add_shape(1, emu(l), emu(t), emu(w), emu(h))
    shape.line.width = Pt(0)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = bw
    return shape

def add_tb(slide, text, l, t, w, h, sz=Pt(11), bold=False, color=DARK,
           align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(emu(l), emu(t), emu(w), emu(h))
    txb.text_frame.word_wrap = wrap
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = sz
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_para(tf, text, sz=Pt(11), bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        pPr = p._p.get_or_add_pPr()
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', str(space_before * 100))
    run = p.add_run()
    run.text = text
    run.font.size = sz
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def remove_pictures(slide):
    """Remove all picture shapes from a slide."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    to_del = [s for s in slide.shapes if s.shape_type == 13]  # 13 = PICTURE
    sp_tree = slide.shapes._spTree
    for s in to_del:
        sp_tree.remove(s._element)
    return len(to_del)

def remove_shapes_by_name(slide, names):
    sp_tree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        if shape.name in names:
            sp_tree.remove(shape._element)
            removed += 1
    return removed

def remove_shapes_in_xrange(slide, x_min_in, x_max_in, exclude_names=None):
    """Remove all shapes whose LEFT position is within the x range (inches)."""
    sp_tree = slide.shapes._spTree
    removed = 0
    for shape in list(slide.shapes):
        l_in = shape.left / 914400 if shape.left else 0
        if x_min_in <= l_in <= x_max_in:
            if exclude_names and shape.name in exclude_names:
                continue
            sp_tree.remove(shape._element)
            removed += 1
    return removed

def find_shape_by_text(slide, fragment):
    for s in slide.shapes:
        if s.has_text_frame and fragment.lower() in s.text_frame.text.lower():
            return s
    return None

def find_shapes_by_text(slide, fragment):
    return [s for s in slide.shapes
            if s.has_text_frame and fragment.lower() in s.text_frame.text.lower()]

def set_para_text(shape, para_idx, text, sz=None, bold=None, color=None):
    tf = shape.text_frame
    if para_idx >= len(tf.paragraphs):
        return
    p = tf.paragraphs[para_idx]
    for r in p.runs:
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = text
    if sz: run.font.size = sz
    if bold is not None: run.font.bold = bold
    if color: run.font.color.rgb = color

def rewrite_textbox(shape, lines, sz=Pt(11), bold=False, color=DARK, first_bold=False, first_sz=None):
    """Replace all text in a textbox with given lines list."""
    tf = shape.text_frame
    txBody = tf._txBody
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)
    for i, line in enumerate(lines):
        is_first = (i == 0)
        p_xml = f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-US" dirty="0"/><a:t>{html.escape(line)}</a:t></a:r></a:p>'
        p_elem = etree.fromstring(p_xml)
        # Set font props via rPr
        rPr = p_elem.find(f'{{{NS}}}r/{{{NS}}}rPr')
        line_sz = first_sz if (is_first and first_sz) else sz
        rPr.set('sz', str(int(line_sz.pt * 100)))
        rPr.set('b', '1' if (bold or (is_first and first_bold)) else '0')
        # color
        c = color
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgb.set('val', f'{c.rgb:06X}' if hasattr(c, 'rgb') else str(c).upper())
        txBody.append(p_elem)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Remove pie chart, restructure to 3 horizontal rows
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 2 ──")
    sl2 = prs.slides[1]

    # 1. Remove the pie chart picture
    n = remove_pictures(sl2)
    print(f"  Removed {n} picture(s)")

    # 2. Remove ALL existing shapes except the header title and subtitle
    #    Keep: shapes with text "The Business Problem" (title) or subtitle text
    #    Delete: all column containers (Rectangles 8-21) and content text boxes
    keep_texts = {"the business problem", "why does predicting loyalty"}
    sp_tree = sl2.shapes._spTree
    to_remove = []
    for shape in list(sl2.shapes):
        if shape.name in ('Rectangle 1', 'Rectangle 2'):  # bg rects — keep
            continue
        txt = shape.text_frame.text.strip().lower() if shape.has_text_frame else ''
        if any(k in txt for k in keep_texts):
            continue
        to_remove.append(shape._element)
    for el in to_remove:
        try: sp_tree.remove(el)
        except: pass
    print(f"  Cleared {len(to_remove)} column shapes")

    # 3. Build 3 horizontal rows
    # Row heights and positions
    rows = [
        {
            'y': 1.25, 'h': 1.65,
            'label': 'WHAT IS SPENDY?',
            'content': [
                'Buy Now, Pay Later (BNPL) platform operating across Japan.',
                'Revenue: merchant transaction fees + consumer instalment interest.',
                'Dataset: 759K transactions, 136.8K users, 711 shops  (Jun – Dec 2020).',
            ]
        },
        {
            'y': 3.05, 'h': 1.75,
            'label': 'THE BRAND LOYALIST PROBLEM',
            'content': [
                '75% of users bought from exactly ONE merchant — they are loyal to the shop, not to Spendy.',
                'If that merchant stops accepting Spendy, the customer disappears entirely.',
                'Raw purchase counts reward these single-shop users — a loyalty programme built on that metric wastes budget.',
            ]
        },
        {
            'y': 4.95, 'h': 1.6,
            'label': 'WHY PREDICT EARLY?',
            'content': [
                'At Month 2 a user\'s habits are still forming — an incentive now is far cheaper than re-acquisition later.',
                'Not all new users have the same potential: the model separates high-value targets from the rest.',
                'Proactive segmentation at Month 2 → differentiated campaigns → better ROI per marketing pound.',
            ]
        },
    ]

    for row in rows:
        y, h = row['y'], row['h']
        # Left accent bar (navy, 0.15" wide)
        add_rect(sl2, 0.3, y + 0.08, 0.12, h - 0.16, fill=NAVY)
        # Row background (very light grey)
        add_rect(sl2, 0.45, y, 12.6, h, fill=LGREY, border=MGREY, bw=Pt(0.5))
        # Label text (top-left of row)
        lbl = sl2.shapes.add_textbox(emu(0.55), emu(y + 0.08), emu(3.5), emu(0.35))
        lbl.text_frame.word_wrap = False
        p = lbl.text_frame.paragraphs[0]
        r = p.add_run(); r.text = row['label']
        r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = NAVY
        # Content text box
        ctb = sl2.shapes.add_textbox(emu(0.55), emu(y + 0.42), emu(12.3), emu(h - 0.5))
        ctb.text_frame.word_wrap = True
        for i, line in enumerate(row['content']):
            if i == 0:
                p = ctb.text_frame.paragraphs[0]
            else:
                p = ctb.text_frame.add_paragraph()
            p.space_before = Pt(2) if i > 0 else Pt(0)
            r = p.add_run(); r.text = '–  ' + line
            r.font.size = Pt(10.5); r.font.color.rgb = DARK

    print("  Built 3 horizontal rows")
except Exception as e:
    print(f"  ERROR in Slide 2: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Remove hand-crafted chart shapes, add EDA images
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 3 ──")
    sl3 = prs.slides[2]

    chart_shape_names = ['Bar1','Bar2','Bar3','Bar4','Bar5','Bar6',
                         'Baseline','MedianLine','PieSliceA','PieSliceB',
                         'Legend1Sq','Legend2Sq']
    n = remove_shapes_by_name(sl3, chart_shape_names)
    print(f"  Removed {n} custom chart shapes")

    # Add frequency distribution chart (left bottom)
    img1 = os.path.join(IMGS, 'eda_frequency_distribution.png')
    if os.path.exists(img1):
        sl3.shapes.add_picture(img1, emu(0.3), emu(4.4), emu(6.2), emu(2.8))
        print(f"  Added eda_frequency_distribution.png")
    else:
        print(f"  SKIP: {img1} not found")

    # Add pie chart (right bottom)
    img2 = os.path.join(IMGS, 'loyalty_merchant_vs_platform_pie.png')
    if os.path.exists(img2):
        sl3.shapes.add_picture(img2, emu(6.7), emu(4.4), emu(6.3), emu(2.8))
        print(f"  Added loyalty_merchant_vs_platform_pie.png")
    else:
        print(f"  SKIP: {img2} not found")

    print("  Added 2 EDA images")
except Exception as e:
    print(f"  ERROR in Slide 3: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Fix OR visibility + space out 4 left items evenly
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 4 ──")
    sl4 = prs.slides[3]

    # Fix OR text box — ensure text is visible (navy color, correct size)
    for s in sl4.shapes:
        if s.has_text_frame and s.text_frame.text.strip() == 'OR':
            # Force font color to NAVY and size to 22pt
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(22)
                    r.font.bold = True
                    r.font.color.rgb = NAVY
            # Also check if it's behind something — bring to front by moving element to end of spTree
            sp_tree = sl4.shapes._spTree
            sp_tree.remove(s._element)
            sp_tree.append(s._element)
            # Make box bigger so text isn't clipped
            s.width = emu(0.8)
            s.height = emu(0.5)
            print(f"  Fixed OR text box at ({s.left/914400:.2f}, {s.top/914400:.2f})")
            break

    # Space out the 4 left requirements to fill y=1.3" to y=6.5"
    item_labels = ['①', '②', '③', '④']
    item_groups = []
    for label in item_labels:
        shapes_in_group = [s for s in sl4.shapes
                           if s.has_text_frame and label in s.text_frame.text]
        if shapes_in_group:
            anchor = min(shapes_in_group, key=lambda s: s.top)
            item_groups.append((label, anchor.top / 914400))

    print(f"  Found {len(item_groups)} item groups: {[(l, round(t,2)) for l,t in item_groups]}")

    target_tops = [1.3, 2.68, 4.06, 5.44]

    for gi, (label, old_top) in enumerate(item_groups):
        new_top = target_tops[gi]
        delta = emu(new_top) - emu(old_top)
        if abs(delta) < 10000:
            continue
        for s in sl4.shapes:
            if not s.has_text_frame:
                continue
            s_top_in = s.top / 914400
            if abs(s_top_in - old_top) < 0.85 and s.left / 914400 < 5.5:
                s.top = s.top + delta
        print(f"  Group {label}: moved from y={old_top:.2f}\" to y={new_top:.2f}\" (delta={delta/914400:.2f}\")")
except Exception as e:
    print(f"  ERROR in Slide 4: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — Update 60-day reasoning + Why It Predicts column
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 5 ──")
    sl5 = prs.slides[4]

    why60 = find_shape_by_text(sl5, '60 days')
    if why60:
        rationale_box = None
        for s in sl5.shapes:
            if s.has_text_frame and s.left/914400 > 2.5 and abs(s.top/914400 - 1.3) < 0.3:
                txt = s.text_frame.text
                if 'leakage' in txt.lower() or 'habit' in txt.lower() or 'only' in txt.lower():
                    rationale_box = s
                    break
        if rationale_box:
            new_lines = [
                'After 2 months, marketing wants to act — this is the earliest point a new user has enough transaction history to score.',
                'Only data from days 0–60 is available at prediction time; using anything beyond this would be data leakage.',
                'The 60-day window also allows a month-over-month spend trend (Month 1 → Month 2) which a 30-day window cannot provide.',
            ]
            rewrite_textbox(rationale_box, new_lines, sz=Pt(10.5), color=DARK)
            print("  Updated WHY 60 DAYS reasoning")
        else:
            print("  WHY 60 DAYS rationale box not found by position")
    else:
        print("  WHY 60 DAYS shape not found")

    why_updates = {
        'high early frequency': 'Users above the 2/month threshold in their first 60 days are demonstrating a habit, not an experiment — early frequency is the strongest predictor of sustained frequency.',
        'already at target spend': 'Reaching ¥10K/month spend within 60 days signals the user is already committing financially — they are not testing, they are buying.',
        'cross-merchant behaviour': 'Visiting a second merchant in 60 days means the user actively chose Spendy in a new context — early breadth is the clearest signal of platform loyalty vs merchant loyalty.',
        'captures early engagement': 'A user who has gone quiet before day 60 is unlikely to sustain loyalty — recency flags early drop-off before it becomes churn.',
        'rising spend': 'Growing spend from Month 1 to Month 2 shows momentum — a user approaching the ¥10K threshold is on a trajectory consistent with loyalty.',
    }

    for fragment, new_text in why_updates.items():
        s = find_shape_by_text(sl5, fragment)
        if s:
            rewrite_textbox(s, [new_text], sz=Pt(10), color=DARK)
            print(f"  Updated 'Why it predicts' for: {fragment[:30]}")
        else:
            print(f"  NOT FOUND: '{fragment[:30]}'")
except Exception as e:
    print(f"  ERROR in Slide 5: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — Remove comparison table, keep only RF reasoning
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 6 ──")
    sl6 = prs.slides[5]

    sp_tree6 = sl6.shapes._spTree
    to_del6 = []
    for shape in list(sl6.shapes):
        l_in = shape.left / 914400 if shape.left else 0
        if l_in >= 7.8:
            to_del6.append(shape._element)
    for el in to_del6:
        try: sp_tree6.remove(el)
        except: pass
    print(f"  Removed {len(to_del6)} comparison table shapes")

    for s in sl6.shapes:
        if s.has_text_frame and s.left/914400 < 1.0:
            l_in = s.left/914400
            t_in = s.top/914400
            if t_in > 1.0:
                s.width = emu(12.5)
    print("  Widened RF reason boxes to full width")

    ci_box = add_rect(sl6, 0.3, 6.2, 12.7, 0.9, fill=RGBColor(0xEB,0xF2,0xF8), border=STEEL, bw=Pt(0.75))
    ci_tb = sl6.shapes.add_textbox(emu(0.5), emu(6.22), emu(12.3), emu(0.8))
    ci_tf = ci_tb.text_frame; ci_tf.word_wrap = True
    p = ci_tf.paragraphs[0]
    r = p.add_run(); r.text = 'CLASS IMBALANCE HANDLING:  '
    r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = NAVY
    r2 = p.add_run()
    r2.text = 'Only 3.2% of training users are loyal.  class_weight=\'balanced\' scales misclassification penalties inversely — the model is penalised more for missing a loyal user than for a false positive.  Stratified split preserves the 3.2% ratio in both train and test sets.'
    r2.font.size = Pt(9.5); r2.font.color.rgb = DARK
    print("  Added CLASS IMBALANCE summary bar")
except Exception as e:
    print(f"  ERROR in Slide 6: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — Rearrange: remove added images, add feature importance chart
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 7 ──")
    sl7 = prs.slides[6]

    n7 = remove_pictures(sl7)
    print(f"  Removed {n7} picture(s)")

    wna = find_shape_by_text(sl7, 'WHY NOT ACCURACY')
    if wna:
        sl7.shapes._spTree.remove(wna._element)
        print("  Removed WHY NOT ACCURACY box")
    else:
        print("  WHY NOT ACCURACY box not found")

    sp_tree7 = sl7.shapes._spTree
    to_del7 = []
    for shape in list(sl7.shapes):
        l_in = shape.left/914400 if shape.left else 0
        if l_in >= 8.3:
            to_del7.append(shape._element)
    for el in to_del7:
        try: sp_tree7.remove(el)
        except: pass
    print(f"  Removed {len(to_del7)} classification report shapes")

    fi_img = os.path.join(IMGS, 'model_feature_importance.png')
    if os.path.exists(fi_img):
        sl7.shapes.add_picture(fi_img, emu(8.3), emu(1.25), emu(4.8), emu(3.8))
        print("  Added feature importance chart to right panel")
    else:
        print(f"  SKIP: {fi_img} not found")

    to_del7b = []
    for shape in list(sl7.shapes):
        t_in = shape.top/914400 if shape.top else 0
        if t_in >= 4.8:
            to_del7b.append(shape._element)
    for el in to_del7b:
        try: sp_tree7.remove(el)
        except: pass
    print(f"  Removed {len(to_del7b)} old feature importance text shapes")

    ki_bg = add_rect(sl7, 0.3, 5.05, 12.7, 2.15, fill=RGBColor(0xEB,0xF2,0xF8), border=STEEL, bw=Pt(0.75))
    ki_tb = sl7.shapes.add_textbox(emu(0.5), emu(5.1), emu(12.3), emu(2.0))
    ki_tf = ki_tb.text_frame; ki_tf.word_wrap = True

    p0 = ki_tf.paragraphs[0]
    r0 = p0.add_run(); r0.text = 'KEY INSIGHT — PRECISION vs RECALL'
    r0.font.size = Pt(10); r0.font.bold = True; r0.font.color.rgb = NAVY

    p1 = ki_tf.add_paragraph(); p1.space_before = Pt(3)
    r1 = p1.add_run()
    r1.text = 'Precision 0.34  |  Recall 0.45  |  F1 0.39  (loyal class)     Not Loyal: Precision 0.98  |  Recall 0.96  |  F1 0.97'
    r1.font.size = Pt(10); r1.font.color.rgb = DARK

    p2 = ki_tf.add_paragraph(); p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = 'Recall 0.45 means the model catches 45% of future loyal users at the default 0.5 threshold.  This is the primary improvement priority — lowering the classification threshold or switching to Gradient Boosting would increase recall at a manageable precision cost.  The High segment (score > 0.66) achieves 83.9% observed loyalty, confirming the model\'s confident predictions are highly accurate.'
    r2.font.size = Pt(10); r2.font.color.rgb = DARK
    print("  Added key insight precision/recall box")
except Exception as e:
    print(f"  ERROR in Slide 7: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Remove added image, shorten descriptions, add image lower
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 8 ──")
    sl8 = prs.slides[7]

    n8 = remove_pictures(sl8)
    print(f"  Removed {n8} picture(s)")

    for s in sl8.shapes:
        if not s.has_text_frame: continue
        t_in = s.top/914400
        if 2.9 <= t_in <= 3.2:
            tf = s.text_frame
            txBody = tf._txBody
            paras = txBody.findall(qn('a:p'))
            content_paras = [p for p in paras if any(
                r.find(qn('a:t')) is not None and r.find(qn('a:t')).text
                for r in p.findall(qn('a:r'))
            )]
            if len(content_paras) > 2:
                for p in paras:
                    txBody.remove(p)
                for p in content_paras[:2]:
                    txBody.append(p)
            s.height = emu(1.0)
            print(f"  Shortened description box at y={t_in:.2f}\"")

    seg_img = os.path.join(IMGS, 'model_segment_distribution.png')
    if os.path.exists(seg_img):
        sl8.shapes.add_picture(seg_img, emu(0.5), emu(4.15), emu(12.0), emu(2.1))
        print("  Added segment distribution chart at y=4.15\"")
    else:
        print(f"  SKIP: {seg_img} not found")
except Exception as e:
    print(f"  ERROR in Slide 8: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Remove Before section, remove diagram, rethink structure
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 9 ──")
    sl9 = prs.slides[8]

    n9 = remove_pictures(sl9)
    print(f"  Removed {n9} picture(s)")

    sp_tree9 = sl9.shapes._spTree
    to_del9 = []
    for shape in list(sl9.shapes):
        l_in = shape.left/914400 if shape.left else 0
        t_in = shape.top/914400 if shape.top else 0
        txt = shape.text_frame.text.strip() if shape.has_text_frame else ''
        if shape.name in ('Rectangle 1', 'Rectangle 2'):
            continue
        if t_in < 1.1:
            continue
        if t_in > 5.1:
            continue
        to_del9.append(shape._element)

    for el in to_del9:
        try: sp_tree9.remove(el)
        except: pass
    print(f"  Removed {len(to_del9)} before/after/arrow shapes")

    subtitle = find_shape_by_text(sl9, 'reactive reporting')
    if subtitle:
        rewrite_textbox(subtitle,
            ['What this analysis delivers — a rigorous loyalty definition and an early-intervention model'],
            sz=Pt(13), color=DARK)
        print("  Updated subtitle")
    else:
        print("  Subtitle 'reactive reporting' not found")

    outcomes = [
        ('LOYALTY DEFINED', '4.6%', 'Platform loyalty rate\nConservative by design\nSolves the Brand Loyalist Problem'),
        ('EARLY PREDICTION', 'AUC 0.912', 'Scores new users at Month 2\nRandom Forest on 5 features\n60-day feature window'),
        ('ACTIONABLE SEGMENTS', '83.9%', 'High segment observed loyalty\n3 tiers with distinct playbooks\nMonth 2 is the intervention window'),
        ('PRODUCTION READY', 'Month 2', 'Monthly batch pipeline\nQuarterly retraining with AUC gate\nA/B testing framework for validation'),
    ]

    card_w = 2.9
    for ci, (title, kpi, desc) in enumerate(outcomes):
        x = 0.35 + ci * (card_w + 0.25)
        add_rect(sl9, x, 1.3, card_w, 3.55, fill=LGREY, border=MGREY, bw=Pt(0.5))
        add_rect(sl9, x, 1.3, card_w, 0.1, fill=NAVY)
        t = sl9.shapes.add_textbox(emu(x+0.1), emu(1.42), emu(card_w-0.2), emu(0.35))
        t.text_frame.word_wrap = False
        r = t.text_frame.paragraphs[0].add_run(); r.text = title
        r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = NAVY
        kpi_tb = sl9.shapes.add_textbox(emu(x+0.1), emu(1.8), emu(card_w-0.2), emu(0.7))
        kpi_tb.text_frame.word_wrap = False
        rk = kpi_tb.text_frame.paragraphs[0].add_run(); rk.text = kpi
        rk.font.size = Pt(26); rk.font.bold = True; rk.font.color.rgb = STEEL
        desc_tb = sl9.shapes.add_textbox(emu(x+0.1), emu(2.55), emu(card_w-0.2), emu(1.2))
        desc_tb.text_frame.word_wrap = True
        for i, line in enumerate(desc.split('\n')):
            p = desc_tb.text_frame.paragraphs[0] if i==0 else desc_tb.text_frame.add_paragraph()
            r = p.add_run(); r.text = '·  ' + line
            r.font.size = Pt(9.5); r.font.color.rgb = DARK

    print("  Built 4 outcome cards")

    auc_box = find_shape_by_text(sl9, '0.912')
    if auc_box:
        rewrite_textbox(auc_box, ['3 Segments'], sz=Pt(32), bold=True, color=STEEL)
        print("  Changed 0.912 KPI to '3 Segments'")
    else:
        print("  0.912 KPI box not found")

    for s in sl9.shapes:
        if s.has_text_frame and 'model' in s.text_frame.text.lower() and 'ranking' in s.text_frame.text.lower():
            rewrite_textbox(s, ['Actionable customer tiers', 'Low / Medium / High'], sz=Pt(11), color=DARK)
            print("  Updated KPI label")
            break
except Exception as e:
    print(f"  ERROR in Slide 9: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Close whitespace in limitation cards
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 10 ──")
    sl10 = prs.slides[9]

    for s in sl10.shapes:
        if s.has_text_frame and s.text_frame.text.strip().upper() == 'HOW TO FIX':
            old_t = s.top / 914400
            s.top = emu(3.1)
            print(f"  Moved HOW TO FIX from y={old_t:.2f}\" to y=3.1\"")

    for s in sl10.shapes:
        if s.has_text_frame:
            t_in = s.top/914400
            txt = s.text_frame.text.strip()
            if 4.5 <= t_in <= 5.0 and len(txt) > 5:
                s.top = emu(3.45)
                print(f"  Moved fix content from y={t_in:.2f}\" to y=3.45\"")
except Exception as e:
    print(f"  ERROR in Slide 10: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Rethink conclusion: DEFINE / PREDICT / ACT 3-column
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 11 ──")
    sl11 = prs.slides[10]

    sp_tree11 = sl11.shapes._spTree
    to_del11 = []
    for shape in list(sl11.shapes):
        if shape.name in ('Rectangle 1', 'Rectangle 5'):
            continue
        txt = shape.text_frame.text.strip() if shape.has_text_frame else ''
        if txt.lower() == 'conclusion':
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(28)
                    r.font.color.rgb = NAVY
            continue
        to_del11.append(shape._element)

    for el in to_del11:
        try: sp_tree11.remove(el)
        except: pass
    print(f"  Cleared {len(to_del11)} old conclusion shapes")

    add_tb(sl11, 'The three things this analysis delivers',
           0.4, 0.85, 12.5, 0.4, sz=Pt(13), color=DARK)

    cols = [
        {
            'title': 'DEFINE',
            'icon_text': '①',
            'color': NAVY,
            'points': [
                ('Loyalty definition', '5-component AND/OR rule — frequency, spend, recency, tenure + platform check'),
                ('4.6% loyalty rate', 'Conservative by design — rigorous enough to attach real rewards to'),
                ('Brand Loyalist Problem', '75% single-shop risk explicitly addressed — merchant loyalty ≠ platform loyalty'),
            ]
        },
        {
            'title': 'PREDICT',
            'icon_text': '②',
            'color': STEEL,
            'points': [
                ('AUC 0.912', 'Random Forest ranks a loyal user above a non-loyal one 91.2% of the time'),
                ('Month 2 scoring', 'First 60 days of transactions are the only input — no leakage, actionable timeline'),
                ('5 features', 'Each maps directly to a loyalty component — spend, frequency, recency, trend, breadth'),
            ]
        },
        {
            'title': 'ACT',
            'icon_text': '③',
            'color': TEAL,
            'points': [
                ('3 actionable segments', 'High (83.9% loyalty) → Reinforce.  Medium (6.3%) → Nudge.  Low (1.0%) → Deprioritise.'),
                ('Differentiated campaigns', 'Each tier has a specific marketing playbook — no more spray-and-pray'),
                ('Validation path', 'A/B experiment design ready — split by model score, measure loyalty rate at 6 months'),
            ]
        },
    ]

    col_w = 3.9
    for ci, col in enumerate(cols):
        x = 0.4 + ci * (col_w + 0.35)
        c = col['color']
        add_rect(sl11, x, 1.35, col_w, 0.55, fill=c)
        h_tb = sl11.shapes.add_textbox(emu(x+0.1), emu(1.38), emu(col_w-0.2), emu(0.48))
        h_tb.text_frame.word_wrap = False
        hp = h_tb.text_frame.paragraphs[0]; hp.alignment = PP_ALIGN.LEFT
        hr = hp.add_run(); hr.text = col['icon_text'] + '  ' + col['title']
        hr.font.size = Pt(14); hr.font.bold = True; hr.font.color.rgb = WHITE
        add_rect(sl11, x, 1.9, col_w, 4.7, fill=LGREY, border=MGREY, bw=Pt(0.5))
        body_tb = sl11.shapes.add_textbox(emu(x+0.15), emu(2.0), emu(col_w-0.3), emu(4.5))
        body_tb.text_frame.word_wrap = True
        first = True
        for (pt_title, pt_body) in col['points']:
            if first:
                p_t = body_tb.text_frame.paragraphs[0]; first = False
            else:
                p_t = body_tb.text_frame.add_paragraph()
                p_t.space_before = Pt(8)
            r_t = p_t.add_run(); r_t.text = pt_title
            r_t.font.size = Pt(11); r_t.font.bold = True; r_t.font.color.rgb = c
            p_b = body_tb.text_frame.add_paragraph()
            r_b = p_b.add_run(); r_b.text = pt_body
            r_b.font.size = Pt(10); r_b.font.color.rgb = DARK

    print("  Built DEFINE / PREDICT / ACT 3-column conclusion")
except Exception as e:
    print(f"  ERROR in Slide 11: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Delete Appendix A (alternatives not in notebook)
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 12 — Delete Appendix A ──")
    xml_slides = prs.slides._sldIdLst
    slide_to_del = prs.slides[11]
    rId = xml_slides[11].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    xml_slides.remove(xml_slides[11])
    try:
        prs.part.drop_rel(rId)
    except Exception as e:
        print(f"  drop_rel note: {e}")
    print("  Deleted Appendix A slide (index 11)")
except Exception as e:
    print(f"  ERROR in Slide 12 deletion: {e}")

# After deletion, Appendix B is now index 11, Appendix C is index 12

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 (was 13) — Appendix A: RF vs Alternatives (confirm + rename)
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 12 (was 13) — Rename to Appendix A ──")
    sl12 = prs.slides[11]
    tag = find_shape_by_text(sl12, 'APPENDIX B')
    if tag:
        rewrite_textbox(tag, ['APPENDIX A'], sz=Pt(8.5), bold=True, color=TEAL)
        print("  Renamed to APPENDIX A")
    else:
        print("  APPENDIX B tag not found on slide 12")
except Exception as e:
    print(f"  ERROR in Slide 12 rename: {e}")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 (was 14) — Appendix B: Rethink pipeline layout
# ═══════════════════════════════════════════════════════════════════
try:
    print("\n── Slide 13 (was 14) — Restructure pipeline detail boxes ──")
    sl13 = prs.slides[12]
    tag13 = find_shape_by_text(sl13, 'APPENDIX C')
    if tag13:
        rewrite_textbox(tag13, ['APPENDIX B'], sz=Pt(8.5), bold=True, color=TEAL)
        print("  Renamed to APPENDIX B")
    else:
        print("  APPENDIX C tag not found on slide 13")

    box_updates = {
        'SCHEDULING': (
            'DATA & FEATURES',
            ['EventBridge triggers monthly pipeline run.',
             'Every new user who crossed day-60 since the last run gets their 60-day feature window computed.',
             'Feature pipeline: Apache Beam or dbt computes the 5 features per user.']
        ),
        'MODEL RETRAINING': (
            'SCORING & SEGMENTATION',
            ['RandomForest.predict_proba() runs on all new feature vectors.',
             'Output: probability score + Low/Medium/High segment label per user.',
             'Results written to CRM/marketing platform within the same day.']
        ),
        'MONITORING': (
            'MONITORING & DRIFT',
            ['KS test on each feature distribution vs baseline — flag if shift > 5%.',
             'Track observed loyalty rate per segment as lagged ground truth (6-month lag).',
             'Alert on score distribution drift before it affects campaign quality.']
        ),
        'THRESHOLD GOVERNANCE': (
            'GOVERNANCE & RETRAINING',
            ['Quarterly model review with AUC >= 0.90 gate before promoting to production.',
             'Segment thresholds (33%/66%) reviewed with marketing & finance each quarter.',
             'A/B experiment results feed back into threshold calibration.']
        ),
    }

    for fragment, (new_title, new_lines) in box_updates.items():
        s = find_shape_by_text(sl13, fragment)
        if s:
            combined = [new_title] + new_lines
            rewrite_textbox(s, combined, sz=Pt(9.5), color=DARK, first_bold=True, first_sz=Pt(10))
            print(f"  Updated box '{fragment}' -> '{new_title}'")
        else:
            print(f"  NOT FOUND: '{fragment}'")
except Exception as e:
    print(f"  ERROR in Slide 13: {e}")

# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════
try:
    prs.save(PPTX)
    print("\n✓ Saved successfully.")
except Exception as e:
    print(f"\n✗ Save failed: {e}")
